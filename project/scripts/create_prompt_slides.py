#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""创建提示词工程培训PPT课件"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs, title, subtitle=None):
    """添加标题幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(46, 117, 182)
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, font_size=20):
    """添加内容幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(46, 117, 182)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(font_size)
        p.space_after = Pt(10)
    
    return slide

def add_comparison_slide(prs, title, left_title, left_content, right_title, right_content):
    """添加左右对比幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(46, 117, 182)
    
    # 左侧标题
    left_title_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.2), Inches(4.5), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(200, 0, 0)
    
    # 左侧内容
    left_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.7), Inches(4.5), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = left_content
    p.font.size = Pt(12)
    
    # 右侧标题
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.5), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 150, 0)
    
    # 右侧内容
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.7), Inches(4.5), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = right_content
    p.font.size = Pt(12)
    
    return slide

def add_code_slide(prs, title, code_lines, font_size=11):
    """添加代码幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(46, 117, 182)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.1), Inches(9.4), Inches(5.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
    shape.line.color.rgb = RGBColor(200, 200, 200)
    
    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    tf = code_box.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(code_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = "Courier New"
        p.space_after = Pt(1)
    
    return slide

def create_ppt():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 封面
    add_title_slide(prs, "提示词工程教程", "Prompt Engineering")
    
    # 目录
    add_content_slide(prs, "课程大纲", [
        "一、什么是提示词（Prompt）",
        "二、设计 Prompt：构建清晰明确的提示",
        "三、使用 Prompt 框架",
        "四、优化 Prompt 的四个技巧",
        "五、Prompt 测试与迭代",
        "六、优化案例分享"
    ])
    
    # 第一部分：什么是提示词
    add_title_slide(prs, "一、什么是提示词", "Prompt 基础概念")
    
    add_content_slide(prs, "提示词（Prompt）的定义", [
        "提示词是您输入给大模型（LLM）的文本信息",
        "用于明确地告诉模型您想要解决的问题或完成的任务",
        "是大语言模型理解用户需求并生成相关、准确回答或内容的基础",
        "提示词质量直接影响模型输出效果"
    ])
    
    # 第二部分：设计Prompt
    add_title_slide(prs, "二、设计 Prompt", "构建清晰明确的提示")
    
    add_content_slide(prs, "构建清晰明确的 Prompt", [
        "想象：给同事指派任务时，只说一句话，效果很难达到预期",
        "如果提供：明确的目的、建议的思考方向、执行策略",
        "结果：更有可能以高标准完成任务",
        "同理：任务描述（Prompt）越清晰、具体、没有歧义",
        "LLM 的表现越能符合您的期望"
    ])
    
    add_comparison_slide(prs, "模糊 vs 清晰的 Prompt",
        "模糊的 Prompt",
        "我想推广公司的新产品。我的公司名为阿里云百炼，新产品名为 Zephyr Z9，是一款轻薄便携的手机。帮我创建一条微博帖子。",
        "清晰具体的 Prompt",
        "请为我司阿里云百炼最新推出的Zephyr Z9轻薄便携手机设计一条吸引眼球的微博推广帖。\n\n内容需彰显Zephyr Z9的独特卖点，如极致轻薄设计、高性能配置及用户便利性，同时融入创意元素以提升观众兴趣和互动意愿。\n\n记得提及阿里云百炼品牌声誉，激发受众好奇心，引导他们探索更多产品信息或直接进行购买。\n\n贴文应简洁有力，符合微博平台的风格与字数限制，适宜社交媒体传播。")
    
    # 第三部分：Prompt框架
    add_title_slide(prs, "三、使用 Prompt 框架", "系统化地向 LLM 提供信息")
    
    add_content_slide(prs, "Prompt 框架的六大要素", [
        "背景：介绍与任务紧密相关的背景信息",
        "目的：明确指出您期望LLM完成的具体任务",
        "风格：指定您希望LLM输出的写作风格",
        "语气：定义输出内容应有的语气（正式、诙谐、温馨等）",
        "受众：明确指出内容面向的读者群体",
        "输出：规定输出内容的具体形式（列表、JSON、报告等）"
    ])
    
    add_comparison_slide(prs, "未使用 vs 使用 Prompt 框架",
        "未使用框架",
        "我想推广公司的新产品。我的公司名为阿里云百炼，新产品名为 Zephyr Z9，是一款轻薄便携的手机。帮我创建一条微博帖子，简洁而深具影响力。",
        "使用 Prompt 框架",
        "#背景# 我想为公司的新产品做广告。我公司的名字叫阿里云百炼，产品叫阿里云百炼 Zephyr Z9，是一款轻薄便携的手机。\n\n#目的# 为我创建一个微博帖子（限制：500字），旨在让人们有兴趣点击产品链接购买。\n\n#风格# 遵循黑米等成功公司为类似产品做广告的写作风格。\n\n#语气# 有说服力\n\n#受众# 我公司在微博上的受众通常是年轻一辈人。\n\n#输出# 微博上的帖子，简洁而有影响力。")
    
    add_content_slide(prs, "使用框架的优势", [
        "框架提醒您考虑需求的各个方面",
        "特别是普通Prompt中缺少的：风格、语气、受众",
        "帮助LLM生成更针对目标群体的内容",
        "细节更多，语言表达更加富有张力",
        "可根据任务需求灵活增减结构组成"
    ])
    
    # 第四部分：优化技巧
    add_title_slide(prs, "四、优化 Prompt 的四个技巧", "提升模型表现")
    
    # 技巧一
    add_content_slide(prs, "技巧一：为模型提供输出样例", [
        "在Prompt中提供您期望的输出示例",
        "让LLM模仿所要求的规范、格式、概念、文法、语气",
        "提供样例可以让大模型多次输出的结果更一致",
        "从而稳定模型表现",
        "特别适用于：格式要求严格、风格统一、批量生成场景"
    ])
    
    add_code_slide(prs, "技巧一示例：小红书种草笔记", [
        "#背景#",
        "你很擅长编写小红书种草笔记，喜欢增加丰富的emoji元素。",
        "",
        "#目的#",
        "请生成一篇小红书种草笔记，推广强森吹风机。",
        "吹风机的优点是：体积小、高颜值、风力大、干得快、智能控温不伤发。",
        "",
        "#受众#",
        "喜欢追求时尚的年轻人，尤其是年轻女性",
        "",
        "#输出#",
        "小红书文章格式，充满emoji元素，简洁但内容充实",
        "",
        "#语气与风格#",
        "我亲测过+n种好物+谁适合谁受益 这个秘诀让你的话语超有信服力！",
        "难题出没+揭秘原因+终极解药 这公式助你条理清晰地分享！",
        "独到见解+深度剖析+巧妙推荐 这公式帮你自然流露心声！",
        "亲身经历+成果展示 这公式让你的情感表达鲜活又感人！"
    ])
    
    # 技巧二
    add_content_slide(prs, "技巧二：设定完成任务的步骤", [
        "对于许多复杂任务，提醒LLM如何完成任务非常必要",
        "通过添加任务步骤，引导模型按步骤执行",
        "特别适用于：数学题、逻辑推理、多步骤任务",
        "帮助模型清晰地从提示词中获取任务步骤",
        "基于规定的步骤生成正确结果"
    ])
    
    add_code_slide(prs, "技巧二示例：数学题求解", [
        "#背景#",
        "小明于早上八点整步行出发，每分钟行走50米，走了12分钟后，",
        "小明的父亲发现小明忘记带作业了，于是便骑车去追小明，",
        "已知小明的爸爸每分钟骑行200米，等到追上小明后，",
        "爸爸决定骑车带上小明，小明坐自行车的路程是走路路程的5倍。",
        "",
        "#目的#",
        "小明什么时候到爷爷家？",
        "",
        "#任务步骤#",
        "1. 先计算小明被爸爸追上时的时间和移动的距离。",
        "2. 再计算小明去爷爷家剩余的距离和需要的时间。",
        "3. 最后计算小明到爷爷家的时间。"
    ])
    
    # 技巧三
    add_content_slide(prs, "技巧三：使用分隔符号区分单元", [
        "在构建复杂的Prompt时，采用特定分隔符界定不同内容单元",
        "显著增强LLM对Prompt正确解析的能力",
        "任务复杂度越高，分隔符的作用越明显",
        "推荐分隔符：###、===、>>> 等罕见字符组合",
        "关键：辨识度高，确保模型明确区分界限标识"
    ])
    
    add_comparison_slide(prs, "技巧三示例：影评总结",
        "未使用分隔符",
        "请简短总结以下影评。\n\n曾经意气风发的张志强在生活的重压下，中年失速偏离了原本的生活轨迹，一时意气用事的决定，让他瞬间从家人的小骄傲变成了社会的边角料。然而，他未曾料到，这仅是他中年人生道路上的起点...",
        "使用分隔符",
        "请简短总结以下影评。\n\n###\n曾经意气风发的张志强在生活的重压下，中年失速偏离了原本的生活轨迹...\n###\n然而，他未曾料到，这仅是他中年人生道路上的起点...\n###\n在共度的时光里他深切的体会到了人间的冷暖真情...\n###")
    
    add_content_slide(prs, "使用分隔符的效果", [
        "LLM理解了三段话的逻辑关系",
        "正确识别了分隔符",
        "根据分隔符的段落生成了三段强相关的总结",
        "输出更加结构化、条理清晰"
    ])
    
    # 技巧四
    add_content_slide(prs, "技巧四：引导模型思考", [
        "对于逻辑推理和语境学习的复杂任务",
        "简单技巧可能无法满足需求",
        "引导模型生成推理过程或帮助模型拆解复杂任务",
        "让模型在生成推理结果前生成更多的推理依据",
        "从而提升模型在复杂问题上的表现"
    ])
    
    add_content_slide(prs, "思维链（Chain of Thought，COT）", [
        "一种使用起来较为简单的引导方法",
        "能够显著提高大模型在复杂场景下的推理能力",
        "核心：让模型先输出思考过程，再输出结论",
        "示例：先输出针对各要求的思考判断过程，再输出最终结论",
        "适用于：JSON解析、逻辑判断、复杂推理等场景"
    ])
    
    add_content_slide(prs, "提示链（Prompt Chaining）", [
        "通过多轮对话，引导LLM思考方向",
        "让LLM从简单任务开始",
        "沿着设计好的思考方向逐步完成复杂推理",
        "构建模式更复杂，但模型表现更好，准确率更高",
        "非常适合逻辑复杂但能按固定模式拆解的困难任务"
    ])
    
    add_content_slide(prs, "其他引导思考的方法", [
        "思维树（Tree of Thoughts, ToT）",
        "Boosting of Thoughts",
        "根据任务特点选择合适的方法",
        "复杂任务可组合使用多种方法"
    ])
    
    # 第五部分：测试与迭代
    add_title_slide(prs, "五、Prompt 测试与迭代", "持续优化")
    
    add_content_slide(prs, "Prompt 优化流程", [
        "1. 设计初始 Prompt",
        "2. 测试模型输出",
        "3. 分析输出问题",
        "4. 调整优化 Prompt",
        "5. 重复测试直到满意",
        "6. 线上环境持续收集反馈"
    ])
    
    add_content_slide(prs, "持续优化的重要性", [
        "生成最优 Prompt 是高度实验性的过程",
        "需要不断尝试和调整各种方法",
        "用户反馈和修正是获取最佳输出的关键",
        "即使优化完成后，仍需持续接收反馈",
        "根据反馈作出相应调整，使模型更好满足用户需求"
    ])
    
    # 第六部分：案例
    add_title_slide(prs, "六、优化案例分享", "跨国公司 AI 助手")
    
    add_content_slide(prs, "案例背景", [
        "问题：qwen-turbo 无法稳定地用英文回答英文问题",
        "场景：跨国公司 HR AI 助手",
        "任务：解答公司政策、考勤制度、年假等问题",
        "挑战：需要稳定输出多语言内容"
    ])
    
    add_content_slide(prs, "优化要点", [
        "1. 将语言替换为语种",
        "   语言有歧义（书面/口头 vs 英语/法语）",
        "   语种没有这个歧义",
        "2. 使用 Prompt 框架重新排版",
        "   原结构较为松散，限制部分内容过于冗余",
        "3. 使用分隔符标记重要内容块",
        "   documents 部分作为独立内容块"
    ])
    
    add_code_slide(prs, "优化后的 Prompt 结构", [
        "#背景#",
        "你是一位跨国公司的高效的HR AI助手，专门负责解答公司内部关于公司政策解析、考勤答疑、年假管理咨询的问题。",
        "以下为公司政策文档：",
        "======",
        "${documents}",
        "======",
        "",
        "#目的#",
        "用户的问题仅限于公司政策解析、考勤答疑、年假管理咨询三类范畴。",
        "当问题在范畴内但知识库未涵盖或不明确时，指引用户联系人力资源部门。",
        "",
        "#多语言要求#",
        "- 如果提出的问题不是中文，用问题的中文含义去检索知识库。",
        "- 知识库检索出来的内容在输出时也要转换为问题的语种。",
        "",
        "#输出#",
        "1. 仅使用标准 ASCII 字符集输出回答。",
        "2. 输出的内容语种为用户输入的语种。"
    ])
    
    # 总结
    add_content_slide(prs, "课程总结", [
        "提示词是与大模型沟通的关键桥梁",
        "清晰具体的 Prompt 是高质量输出的基础",
        "Prompt 框架帮助系统化构建提示词",
        "四个优化技巧：样例、步骤、分隔符、引导思考",
        "持续测试迭代是优化的必经之路"
    ])
    
    # 结束语
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "掌握提示词工程"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(46, 117, 182)
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "让AI成为您的得力助手"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "感谢学习！"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER
    
    # 保存
    output_path = 'docs/training/ai-prompt-training/prompt-engineering-slides.pptx'
    prs.save(output_path)
    print(f'提示词工程培训PPT已创建：{output_path}')

if __name__ == '__main__':
    create_ppt()
