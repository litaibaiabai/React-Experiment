#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""创建AI提示词培训PPT课件"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs, title, subtitle=None):
    """添加标题幻灯片"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(46, 117, 182)
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        # 副标题
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets):
    """添加内容幻灯片"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(46, 117, 182)
    
    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(20)
        p.space_after = Pt(12)
    
    return slide

def add_comparison_slide(prs, title, wrong_text, right_text):
    """添加对比幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(46, 117, 182)
    
    # 错误示范
    wrong_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(2))
    tf = wrong_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "❌ 错误示范"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(200, 0, 0)
    p = tf.add_paragraph()
    p.text = wrong_text
    p.font.size = Pt(18)
    p.space_before = Pt(8)
    
    # 正确示范
    right_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(2.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "✅ 正确示范"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 150, 0)
    p = tf.add_paragraph()
    p.text = right_text
    p.font.size = Pt(16)
    p.space_before = Pt(8)
    
    return slide

def add_code_slide(prs, title, code_lines):
    """添加代码幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(46, 117, 182)
    
    # 代码框背景
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.2), Inches(9.4), Inches(5.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
    shape.line.color.rgb = RGBColor(200, 200, 200)
    
    # 代码内容
    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5))
    tf = code_box.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(code_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.name = "Courier New"
        p.space_after = Pt(2)
    
    return slide

def add_table_slide(prs, title, headers, data):
    """添加表格幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(46, 117, 182)
    
    # 表格
    rows = len(data) + 1
    cols = len(headers)
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.3), Inches(9), Inches(0.5 * rows)).table
    
    # 表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(232, 244, 248)
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(16)
    
    # 数据
    for i, row_data in enumerate(data):
        for j, text in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = text
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
    
    return slide

def create_ppt():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 1. 封面
    add_title_slide(prs, "AI助手进课堂", "小学教师提示词入门")
    
    # 2. 破冰问题
    add_content_slide(prs, "破冰问题", [
        "备一节课需要多长时间？",
        "1小时？2小时？更久？",
        "如果30秒能生成教案框架呢？"
    ])
    
    # 3. 神奇时刻
    add_content_slide(prs, "神奇时刻", [
        "现场演示：30秒生成教案",
        "AI工具：Kimi/文心一言/豆包",
        "让我们一起见证AI的力量"
    ])
    
    # 4-6. AI是什么
    add_content_slide(prs, "AI是什么？", [
        "AI = 超级助手",
        "需要你给它清晰的指令",
        "提示词就是给助手的「工作说明」"
    ])
    
    add_content_slide(prs, "AI不会取代老师", [
        "AI是工具，不是替代者",
        "会用AI的老师效率会高10倍",
        "今天你就能学会！"
    ])
    
    add_content_slide(prs, "今日学习目标", [
        "掌握5个核心提示词技巧",
        "获得3个可直接使用的模板",
        "回去就能用AI辅助备课"
    ])
    
    # 案例1：生成教案
    add_comparison_slide(prs, "案例1：生成教案框架",
        "「帮我写一个《荷花》的教案」",
        "「你是一位有20年经验的小学语文教师，擅长情境教学。\n请帮我设计一节三年级语文《荷花》的教案框架：\n- 课题：《荷花》\n- 课时：2课时\n- 请提供：教学目标、教学重难点、教学过程、板书设计」")
    
    add_content_slide(prs, "技巧1：角色设定", [
        "告诉AI它是谁",
        "它会以这个身份来思考和回答",
        "示例：「你是一位有20年经验的小学语文教师」"
    ])
    
    add_content_slide(prs, "技巧2：任务描述", [
        "明确告诉AI你要什么",
        "越具体越好",
        "关键要素：学科、年级、课题、课时、要求"
    ])
    
    add_code_slide(prs, "模板1：教案生成", [
        "你是一位有20年经验的小学【学科】教师，擅长【教学特色】。",
        "",
        "请帮我设计一节【年级】【学科】课的教案框架：",
        "- 课题：【课题名称】",
        "- 课时：【X】课时",
        "- 教学目标需要包含：知识与技能、过程与方法、情感态度价值观",
        "- 请提供：教学目标、教学重难点、教学过程、板书设计"
    ])
    
    # 案例2：设计活动
    add_comparison_slide(prs, "案例2：设计课堂活动",
        "「设计一个分数教学的活动」",
        "「你是一位擅长游戏化教学的小学数学教师。\n请帮我设计一个四年级《分数的初步认识》的课堂活动：\n- 活动时长：15分钟\n- 参考风格：像「小组竞赛」那样\n- 请提供：活动名称、规则、材料、引导语」")
    
    add_content_slide(prs, "技巧3：提供示例", [
        "给AI看一个你喜欢的范例",
        "它会模仿这个风格",
        "示例：「参考风格：像抢答游戏那样」"
    ])
    
    add_content_slide(prs, "技巧4：迭代优化", [
        "第一次不满意？追问修改！",
        "「请把活动时间缩短到10分钟」",
        "「请增加一个互动环节」"
    ])
    
    add_code_slide(prs, "模板2：活动设计", [
        "你是一位擅长游戏化教学的小学【学科】教师。",
        "",
        "请帮我设计一个【年级】【课题】的课堂活动：",
        "- 活动时长：【X】分钟",
        "- 活动目标：【具体目标】",
        "- 参考风格：【示例描述】",
        "",
        "请提供：",
        "1. 活动名称  2. 活动规则  3. 所需材料  4. 教师引导语"
    ])
    
    # 案例3：制作课件
    add_comparison_slide(prs, "案例3：制作课件内容",
        "「帮我做一个地球运动的课件」",
        "「你是一位擅长制作教学课件的小学科学教师。\n请帮我设计一个五年级《地球的运动》的PPT课件内容大纲：\n- 课件页数：约10页\n- 请按格式输出：每页标题、核心内容、配图建议、讲解提示」")
    
    add_content_slide(prs, "技巧5：结构化输出要求", [
        "明确告诉AI你要什么格式",
        "「请按以下格式输出每一页...」",
        "输出更规范，直接可用"
    ])
    
    add_code_slide(prs, "模板3：课件内容", [
        "你是一位擅长制作教学课件的小学【学科】教师。",
        "",
        "请帮我设计一个【年级】【课题】的PPT课件内容大纲：",
        "- 课件页数：约【X】页",
        "",
        "请按以下格式输出每一页：",
        "【第X页】标题：______",
        "- 核心内容（3-5个要点）",
        "- 配图建议",
        "- 讲解提示"
    ])
    
    # 练习环节
    add_content_slide(prs, "动手练习（20分钟）", [
        "任务：选择一个下周要上的课题",
        "套用模板1生成教案框架",
        "讲师巡视，解答问题",
        "分享成果，互相学习"
    ])
    
    # 总结
    add_table_slide(prs, "5个技巧回顾", 
        ["技巧", "说明", "示例"],
        [
            ["角色设定", "告诉AI它是谁", "「你是一位资深教师」"],
            ["任务描述", "明确你要什么", "「设计三年级《荷花》教案」"],
            ["提供示例", "给AI看范例", "「像抢答游戏那样」"],
            ["迭代优化", "追问修改", "「请把导入改成游戏」"],
            ["结构化输出", "指定格式", "「请按以下格式输出」"]
        ])
    
    add_content_slide(prs, "3个模板汇总", [
        "模板1：教案生成 - 快速生成教案框架",
        "模板2：活动设计 - 设计有趣的课堂活动",
        "模板3：课件内容 - 制作PPT内容大纲"
    ])
    
    add_table_slide(prs, "推荐AI工具",
        ["工具", "特点", "适用场景"],
        [
            ["Kimi", "免费、长文本", "教案、课件"],
            ["文心一言", "中文优化", "各类场景"],
            ["豆包", "新手友好", "入门首选"],
            ["通义千问", "功能全面", "各类场景"]
        ])
    
    add_content_slide(prs, "后续学习建议", [
        "多用：每天尝试用AI完成一个小任务",
        "多改：不满意时修改提示词",
        "多交流：和同事分享好用的提示词"
    ])
    
    # 结束语
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AI不会取代老师"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(46, 117, 182)
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "但会用AI的老师会更轻松"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(46, 117, 182)
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "祝各位老师工作顺利！"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER
    
    # 保存
    output_path = 'docs/training/ai-prompt-training/training-slides.pptx'
    prs.save(output_path)
    print(f'培训PPT课件已创建：{output_path}')

if __name__ == '__main__':
    create_ppt()
